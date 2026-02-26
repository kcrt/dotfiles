#!/usr/bin/env -S uv run
# -*- coding: utf-8 -*-
# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "python-docx",
#     "python-pptx",
#     "openpyxl",
#     "pypdf",
#     "rarfile",
# ]
# ///

"""
exstrings.py - Extract text content from various file types.

Supported file types: docx, pptx, xlsx, pdf, rtf, doc, zip, rar, lzh

Can read from file path or stdin (for use with ripgrep-all).
"""

import argparse
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from pypdf import PdfReader
from pptx import Presentation


def extract_text_from_docx(filepath: Path) -> str:
    """Extract text from a DOCX file."""
    try:
        doc = DocxDocument(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except zipfile.BadZipFile:
        # Encrypted or corrupted file - return empty string silently
        return ""


def extract_text_from_pptx(filepath: Path) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(filepath)
    slides_text = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            slides_text.append(f"----- Slide {idx} -----")
            slides_text.append("\n".join(slide_texts))

    return "\n".join(slides_text)


def extract_text_from_xlsx(filepath: Path) -> str:
    """Extract text from an XLSX file."""
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
        sheets_text = []

        for sheet in wb:
            sheets_text.append(f"----- Sheet: {sheet.title} -----")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(v) if v is not None else "" for v in row)
                if row_text.strip():
                    sheets_text.append(row_text)

        return "\n".join(sheets_text)
    except (zipfile.BadZipFile, InvalidFileException):
        # Encrypted or corrupted file - return empty string silently
        return ""


def extract_text_from_pdf(filepath: Path) -> str:
    """Extract text from a PDF file."""
    reader = PdfReader(filepath)
    text_parts = []

    for page_num, page in enumerate(reader.pages, 1):
        page_text = page.extract_text()
        if page_text.strip():
            text_parts.append(f"----- Page {page_num} -----")
            text_parts.append(page_text)

    return "\n".join(text_parts)


def extract_text_via_textutil(filepath: Path) -> str:
    """Extract text using macOS textutil command."""
    if not shutil.which("textutil"):
        return "textutil not found (macOS only)"

    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(filepath)],
        capture_output=True,
        text=True,
    )
    return result.stdout


def list_zip_contents(filepath: Path) -> str:
    """List contents of a ZIP file."""
    lines = []
    with zipfile.ZipFile(filepath) as zf:
        for info in zf.infolist():
            lines.append(info.filename)
    return "\n".join(lines)


def list_rar_contents(filepath: Path) -> str:
    """List contents of a RAR file using unrar."""
    if not shutil.which("unrar"):
        return "unrar command not found"

    result = subprocess.run(
        ["unrar", "l", str(filepath)],
        capture_output=True,
        text=True,
    )
    return result.stdout


def list_lzh_contents(filepath: Path) -> str:
    """List contents of an LZH file using lha."""
    for cmd in ["lha", "lha"]:
        if shutil.which(cmd):
            result = subprocess.run(
                [cmd, "l", str(filepath)],
                capture_output=True,
                text=True,
            )
            return result.stdout
    return "lha command not found"


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Extract text content from various file types."
    )
    parser.add_argument("file", nargs="?", help="File to extract text from (omit for stdin)")
    parser.add_argument("--type", help="File type (required when reading from stdin)")
    parser.add_argument("--stdin", action="store_true", help="Read from stdin")
    args = parser.parse_args()

    use_stdin = args.stdin or args.file is None

    if use_stdin:
        # Read from stdin, write to temp file
        if not args.type:
            print("Error: --type is required when reading from stdin", file=sys.stderr)
            return 1
        ext = args.type.lower()
        # Read stdin to temp file
        with tempfile.NamedTemporaryFile(mode='wb', delete=False, suffix=f'.{ext}') as f:
            temp_path = Path(f.name)
            f.write(sys.stdin.buffer.read())
        filepath = temp_path
    else:
        filepath = Path(args.file)
        if not filepath.exists():
            print(f"Error: File not found: {filepath}", file=sys.stderr)
            return 1
        ext = filepath.suffix.lstrip(".").lower()

    try:
        handlers = {
            "docx": extract_text_from_docx,
            "pptx": extract_text_from_pptx,
            "xlsx": extract_text_from_xlsx,
            "pdf": extract_text_from_pdf,
            "rtf": extract_text_via_textutil,
            "doc": extract_text_via_textutil,
            "zip": list_zip_contents,
            "rar": list_rar_contents,
            "lzh": list_lzh_contents,
        }

        handler = handlers.get(ext)
        if handler is None:
            print(f"Error: Unsupported file type: .{ext}", file=sys.stderr)
            return 1

        result = handler(filepath)
        if result:
            print(result)
    finally:
        # Clean up temp file if we created one
        if use_stdin:
            filepath.unlink()

    return 0


if __name__ == "__main__":
    sys.exit(main())
