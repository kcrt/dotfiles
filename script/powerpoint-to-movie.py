#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# needs: pip3 install python-pptx gtts moviepy pillow pdf2image
# needs(macOS):

import shutil
import subprocess
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from moviepy.editor import *
from gtts import gTTS
from pdf2image import convert_from_path
import os
import platform
import argparse
import re


def parse_args():
    parser = argparse.ArgumentParser(
        description='Convert a PowerPoint presentation to a video.')
    parser.add_argument('pptx_path', type=str,
                        help='The path to the PowerPoint file.')
    parser.add_argument('output_path', type=str,
                        help='The path to save the output video file.')
    parser.add_argument('--lang', type=str, default='ja',
                        help='The language for the text-to-speech conversion.')
    parser.add_argument('--no_audio_duration', type=float, default=5.0,
                        help='The duration of the slide if audio is not available.')
    parser.add_argument('--speed', type=float, default=1.0,
                        help='The speed of the video. 1.0 is the normal speed.')
    parser.add_argument('--padding', type=float, default=1,
                        help='The padding duration for each slide with audio.')
    parser.add_argument('--codec', type=str, default='libx264',
                        help='The codec to use for the output video.')
    parser.add_argument('--pronunciation', type=str,
                        help="The pronunciation csv file. The second column should be the pronunciation of the first column."),
    return parser.parse_args()


def extract_notes(pptx_path: str) -> List[str]:
    """Extracts notes from a PowerPoint presentation.
    Notes are converted accprding to the following rules:
    1. remove lines after "---" in notes
    2. <名称|よみかた> will be converted to "よみかた"

    Args:
        pptx_path (str): The path to the PowerPoint file.

    Returns:
        List[str]: A list of notes extracted from each slide.
    """
    prs = Presentation(pptx_path)
    notes = []

    for slide in prs.slides:
        if slide.notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            # Rule 1: remove lines after "---"
            if "---" in notes_text:
                notes_text = notes_text.split("---")[0]
            # Rule 2: <名称|よみかた> will be converted to "よみかた"
            notes_text = re.sub(r'<[^>]*\|([^>]*)>', r'\1', notes_text)

            notes.append(notes_text)
        else:
            notes.append("")

    return notes


def apply_pronunciation(note: str, pronunciation_dict: dict[str, str]) -> str:
    """Applies pronunciation to a note.

    Args:
        note (str): The note to apply pronunciation.
        pronunciation_dict (dict[str, str]): The dictionary of pronunciation. (expression: pronunciation)

    Returns:
        str: The note with pronunciation applied.
    """
    for key, value in pronunciation_dict.items():
        note = note.replace(key, value)

    return note


def pptx_to_pdf_with_windows_powerpoint(pptx_path: str, output_pdf_path: str) -> None:
    """Converts a PowerPoint file to a PDF.

    Args:
        pptx_path (str): The path to the PowerPoint file.
        output_pdf_path (str): The path to save the resulting PDF.
    """
    import comtypes.client
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(pptx_path)
    presentation.SaveAs(output_pdf_path, 32)  # formatType = 32 for ppt to pdf
    presentation.Close()
    powerpoint.Quit()


def pptx_to_pdf_with_libreoffice(pptx_path: str, output_pdf_path: str) -> None:
    """
    Converts a PowerPoint presentation to a PDF file using LibreOffice.
    """
    if platform.system() == "Darwin":
        libreoffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    else:
        libreoffice_path = "libreoffice"
    command = [
        libreoffice_path,
        '--headless',
        '--convert-to',
        'pdf',
        '--outdir',
        '/tmp',
        pptx_path
    ]
    print(" ".join(command))
    subprocess.run(command, check=True)
    converted_pdf_filename = os.path.basename(
        pptx_path).replace('.pptx', '.pdf')
    converted_pdf_path = f'/tmp/{converted_pdf_filename}'

    os.rename(converted_pdf_path, output_pdf_path)


def pdf_to_images(pdf_path: str, output_folder: str) -> List[str]:
    """Converts each page of a PDF file to an image.

    Args:
        pdf_path (str): The path to the PDF file.
        output_folder (str): The folder to save the images.

    Returns:
        List[str]: A list of paths to the saved images.
    """
    pages = convert_from_path(pdf_path)
    image_paths = []

    for i, page in enumerate(pages):
        image_path = f"{output_folder}/slide_{i+1:02d}.png"
        page.save(image_path, 'PNG')
        image_paths.append(image_path)

    return image_paths


def text_to_speech(text: str, output_path: str, lang: str = 'en') -> str:
    """Converts text to speech and saves it as an MP3 file.

    Args:
        text (str): The text to be converted to speech.
        output_path (str): The path to save the MP3 file.
        lang (str, optional): The language for the text-to-speech conversion. Defaults to 'en'.

    Returns:
        str: The path to the saved MP3 file.
    """
    tts = gTTS(text=text, lang=lang)
    tts.save(output_path)
    return output_path


def create_video(files_path: str, output_path: str, no_audio_duration: float = 5.0, speed: float = 1.0, padding: float = 1, codec="libx264") -> None:
    """Creates a video from images and audio files.

    Args:
        files_path (str): The path to the image and audio file. Image files should be named like slide_01.png. Audio files should be named like note_01.mp3.
        output_path (str): The path to save the output video file.
        no_audio_duration (float): The duration of the slide if audio is not available.
        padding (float): The padding duration for each slide with audio.
        speed (float): The speed of the video. 1.0 is the normal speed.
    """
    # Get sorted list of image files
    image_files = sorted(
        [f for f in os.listdir(files_path) if f.endswith('.png')])

    # Create list of clips
    clips = []
    for image_file in image_files:
        image_path = os.path.join(files_path, image_file)
        slide_clip = ImageClip(image_path)
        index = int(image_file.split('_')[1].split('.')[0])

        # Check if corresponding audio file exists
        audio_file = f'note_{index:02d}.mp3'
        audio_path = os.path.join(files_path, audio_file)

        if os.path.exists(audio_path):
            audio_clip = AudioFileClip(audio_path)
            slide_clip = slide_clip.set_audio(
                audio_clip).set_duration(audio_clip.duration + padding)
        else:
            slide_clip = slide_clip.set_duration(no_audio_duration)

        clips.append(slide_clip)

    # Concatenate all clips
    video = concatenate_videoclips(clips, method="compose")

    # Adjust video speed
    final_video = video.fx(vfx.speedx, speed)

    # Write the video file
    final_video.write_videofile(output_path, codec=codec, fps=10)


def main():

    args = parse_args()

    # if /tmp/slides.pdf and /tmp/slides exist, remove them
    if os.path.exists("/tmp/slides.pdf"):
        os.remove("/tmp/slides.pdf")
    if os.path.exists("/tmp/slides"):
        shutil.rmtree("/tmp/slides")

    # check if file exists
    if not os.path.exists(args.pptx_path):
        print(f"File not found: {args.pptx_path}")
        return

    print(f"Converting {args.pptx_path} to PDF...")
    if platform.system() == "Windows":
        pptx_to_pdf_with_windows_powerpoint(args.pptx_path, "/tmp/slides.pdf")
    else:
        pptx_to_pdf_with_libreoffice(args.pptx_path, "/tmp/slides.pdf")
    os.mkdir("/tmp/slides")
    print("Converting PDF to images...")
    image_paths = pdf_to_images("/tmp/slides.pdf", "/tmp/slides")

    print("Extracting notes...")
    notes = extract_notes(args.pptx_path)

    if args.pronunciation:
        print("Applying pronunciation...")
        pronunciation_dict = {}
        with open(args.pronunciation, "r") as f:
            # skip first line
            f.readline()
            for line in f:
                # remove leading and trailing whitespaces
                expression, pronunciation = [
                    x.strip() for x in line.split(",")]
                pronunciation_dict[expression] = pronunciation
        notes = [apply_pronunciation(note, pronunciation_dict)
                 for note in notes]

    print("Converting notes to speech...")
    for i, note in enumerate(notes):
        print(f"Slide #{i+1}: {note}")
        if note != "":
            text_to_speech(
                note, f"/tmp/slides/note_{i+1:02d}.mp3", lang=args.lang)

    print("Creating video...")
    create_video("/tmp/slides", args.output_path,
                 args.no_audio_duration, args.speed, args.padding, args.codec)

    os.remove("/tmp/slides.pdf")
    shutil.rmtree("/tmp/slides")


if __name__ == '__main__':
    main()
